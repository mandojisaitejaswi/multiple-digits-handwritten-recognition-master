from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Decoding Handwriting: Building a Smart Digit Recognition System"
subtitle.text = "A Convolutional Neural Network Approach"

# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "Background:\n"
    " - Millions of handwritten digits are processed daily in banking and postal services.\n"
    "Why It Matters:\n"
    " - Accurate digit recognition is essential for automation.\n"
    "Problem Statement:\n"
    " - From scribbles to precision: Making machines understand handwriting!"
)

# Slide 3: Dataset
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Dataset Overview"
content.text = (
    " - Source: MNIST Dataset\n"
    " - Features: Grayscale, 28x28 resolution\n"
    " - Distribution: Balanced across digits 0-9\n"
    " - Visualized sample data"
)

# Slide 4: Dataset Cleaning
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Dataset Cleaning"
content.text = (
    " - Challenges: Noisy data, overlapping digits\n"
    " - Techniques:\n"
    "   - Noise reduction (Gaussian Blurring)\n"
    "   - Contour extraction with OpenCV\n"
    "   - Outlier detection using PCA"
)

# Slide 5: Model Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Model Architecture"
content.text = (
    " - Convolutional Neural Network (CNN)\n"
    " - Layers:\n"
    "   - Convolution + ReLU (Feature extraction)\n"
    "   - MaxPooling (Dimensionality reduction)\n"
    "   - Fully Connected for Classification\n"
    " - Hyperparameters:\n"
    "   - Optimizer: Adam\n"
    "   - Learning Rate: 0.001"
)

# Slide 6: Training and Validation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Training and Validation"
content.text = (
    " - Data Splits:\n"
    "   - Train: 80%, Validation: 10%, Test: 10%\n"
    " - Techniques:\n"
    "   - Data Augmentation\n"
    "   - Early Stopping\n"
    " - Metrics:\n"
    "   - Accuracy, Precision, Recall, F1-Score"
)

# Slide 7: Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Results"
content.text = (
    " - Training Accuracy: 98%\n"
    " - Validation Accuracy: 95%\n"
    " - Visualizations:\n"
    "   - Confusion Matrix\n"
    "   - Loss and Accuracy Curves"
)

# Slide 8: Experimentation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Experimentation"
content.text = (
    " - Model Comparisons:\n"
    "   - CNN vs. SVM vs. Random Forest\n"
    " - Ablation Studies:\n"
    "   - Impact of Dropout Layers\n"
    "   - Learning Rate Tuning"
)

# Slide 9: Real-Time Application (GUI)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Real-Time Application (GUI)"
content.text = (
    " - Features:\n"
    "   - Real-time digit recognition using Tkinter\n"
    "   - Confidence Scores displayed\n"
    " - Workflow:\n"
    "   - Input → Preprocessing → Prediction → Output"
)

# Slide 10: Future Work
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Future Work"
content.text = (
    " - Extend capabilities to multilingual text recognition\n"
    " - Deploy the system as a web application\n"
    " - Experiment with Transformer-based models for image recognition"
)

# Slide 11: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    " - Successfully built a robust ML pipeline for digit recognition.\n"
    " - Demonstrated real-time recognition with GUI.\n"
    " - Future potential in OCR systems and mobile deployment."
)

# Save the presentation
prs.save("ML_Project_Presentation.pptx")
print("Presentation created and saved as 'ML_Project_Presentation.pptx'")
